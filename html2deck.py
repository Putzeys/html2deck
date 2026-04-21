#!/usr/bin/env python3
"""
html2deck - Convert any HTML presentation to PDF, PPTX, and/or MP4.

PDF/PPTX: screenshot per slide (static images).
MP4: live browser recording with all CSS animations and transitions.

Usage:
    html2deck slides.html                     # PDF + PPTX
    html2deck slides.html -f mp4              # MP4 video with effects
    html2deck slides.html -f mp4 --duration 5 # 5 seconds per slide
    html2deck slides.html -f all              # PDF + PPTX + MP4
"""

from __future__ import annotations

import argparse
import asyncio
import hashlib
import http.server
import io
import subprocess
import sys
import tempfile
import threading
from pathlib import Path

try:
    from PIL import Image
    from playwright.async_api import async_playwright
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as exc:
    name = str(exc).split("'")[1] if "'" in str(exc) else str(exc)
    print(f"Missing: {name}", file=sys.stderr)
    print("pip install playwright python-pptx Pillow && playwright install chromium", file=sys.stderr)
    sys.exit(1)


# ── JS snippets ──────────────────────────────────────────────────────

JS_DETECT_FRAMEWORK = """() => {
    if (typeof Reveal !== 'undefined' && typeof Reveal.slide === 'function') {
        const idx = [];
        document.querySelectorAll('.reveal .slides > section').forEach((s, h) => {
            const v = Array.from(s.children).filter(c => c.tagName === 'SECTION');
            v.length > 0 ? v.forEach((_, vi) => idx.push([h, vi])) : idx.push([h, 0]);
        });
        return { fw: 'revealjs', total: idx.length, meta: idx };
    }
    if (typeof slideshow !== 'undefined' && typeof slideshow.gotoSlide === 'function')
        return { fw: 'remarkjs', total: slideshow.getSlideCount(), meta: null };
    if (typeof impress !== 'undefined' && document.querySelectorAll('.step').length > 1)
        return { fw: 'impressjs', total: document.querySelectorAll('.step').length, meta: null };
    const sr = document.querySelector('.shower')
        || (document.documentElement.classList.contains('shower') ? document.documentElement : null);
    if (sr) { const s = sr.querySelectorAll('.slide'); if (s.length > 1) return { fw: 'shower', total: s.length, meta: null }; }
    return null;
}"""

JS_GOTO = """({fw, index, meta}) => {
    if (fw === 'revealjs') Reveal.slide(meta[index][0], meta[index][1]);
    else if (fw === 'remarkjs') slideshow.gotoSlide(index + 1);
    else if (fw === 'impressjs') impress().goto(document.querySelectorAll('.step')[index]);
    else if (fw === 'shower') {
        const s = (document.querySelector('.shower')||document.documentElement).querySelectorAll('.slide');
        location.hash = s[index]?.id ? '#'+s[index].id : '#'+(index+1);
    }
}"""

JS_DOM_COUNT = """() => {
    const vw = window.innerWidth;
    for (const sel of ['.slide-item','.slide','.swiper-slide','.step','[data-slide]','section']) {
        const els = document.querySelectorAll(sel);
        if (els.length >= 2 && els[0].getBoundingClientRect().width >= vw * 0.7)
            return { sel, n: els.length };
    }
    return null;
}"""

JS_FONTS = "() => document.fonts.ready"
JS_ZOOM = """(z) => { document.body.style.zoom = z; }"""
JS_DOC_H = "() => document.documentElement.scrollHeight"
JS_VH = "() => window.innerHeight"

JS_ISOLATE = """({sel, idx}) => {
    const els = document.querySelectorAll(sel);
    const c = els[0]?.parentElement;
    if (c) { ['transition','transform','display','overflow'].forEach(p =>
        c.style.setProperty(p, p==='display'?'block':'none','important')); }
    els.forEach((el, j) => {
        if (j === idx) {
            Object.entries({position:'fixed',top:'0',left:'0',width:'100vw',height:'100vh',
                'z-index':'99999',display:'flex','align-items':'center','justify-content':'center',
                'flex-direction':'column',visibility:'visible',opacity:'1',transform:'none',
                overflow:'hidden','box-sizing':'border-box'}).forEach(([k,v]) =>
                el.style.setProperty(k,v,'important'));
        } else el.style.setProperty('display','none','important');
    });
}"""


# ── HTTP server ──────────────────────────────────────────────────────

class _Q(http.server.SimpleHTTPRequestHandler):
    def log_message(self, *_): pass

def _serve(d: Path):
    class H(_Q):
        def __init__(self, *a, **k): super().__init__(*a, directory=str(d), **k)
    s = http.server.HTTPServer(("127.0.0.1", 0), H)
    threading.Thread(target=s.serve_forever, daemon=True).start()
    return s, s.server_address[1]


# ── Helpers ──────────────────────────────────────────────────────────

def _h(d: bytes) -> str: return hashlib.md5(d).hexdigest()

def _p(cur, tot):
    print(f"\r  Slide {cur}" + (f"/{tot}" if tot else ""), end="", flush=True)


# ── Screenshot capture strategies ────────────────────────────────────

async def cap_framework(page, info, wait):
    imgs = []
    for i in range(info["total"]):
        await page.evaluate(JS_GOTO, {"fw": info["fw"], "index": i, "meta": info["meta"]})
        await page.wait_for_timeout(wait)
        imgs.append(await page.screenshot(type="png"))
        _p(i+1, info["total"])
    print(); return imgs

async def cap_keyboard(page, wait, hint):
    imgs, prev, same = [], None, 0
    for i in range(hint or 300):
        if i > 0:
            await page.keyboard.press("ArrowRight")
            await page.wait_for_timeout(wait)
        png = await page.screenshot(type="png")
        h = _h(png)
        if h == prev:
            same += 1
            if same >= 2: break
            continue
        same = 0; imgs.append(png); prev = h; _p(len(imgs), hint)
    print(); return imgs

async def cap_scroll(page, wait):
    th = await page.evaluate(JS_DOC_H)
    vh = await page.evaluate(JS_VH)
    n = max(1, -(-th // vh)); imgs = []
    for i in range(n):
        await page.evaluate(f"window.scrollTo(0,{i*vh})")
        await page.wait_for_timeout(wait)
        imgs.append(await page.screenshot(type="png")); _p(i+1, n)
    print(); return imgs

async def cap_selector(page, sel, wait):
    c = await page.evaluate(f"document.querySelectorAll('{sel}').length")
    if c == 0: print(f"  No elements match '{sel}'"); return []
    imgs = []
    for i in range(c):
        await page.evaluate(JS_ISOLATE, {"sel": sel, "idx": i})
        await page.wait_for_timeout(wait)
        imgs.append(await page.screenshot(type="png")); _p(i+1, c)
    print(); return imgs


# ── Capture slides (detect + capture) ────────────────────────────────

async def _apply_zoom(page, zoom):
    if zoom != 100:
        await page.evaluate(JS_ZOOM, zoom / 100)
        await page.wait_for_timeout(500)

async def capture_slides(pw, url, args):
    browser = await pw.chromium.launch()
    page = await browser.new_page(
        viewport={"width": args.width, "height": args.height},
        device_scale_factor=args.scale,
    )
    await page.goto(url, wait_until="networkidle", timeout=30_000)
    await page.evaluate(JS_FONTS)
    await page.wait_for_timeout(1200)
    await _apply_zoom(page, args.zoom)

    if args.selector:
        print(f"  Strategy: selector '{args.selector}'")
        imgs = await cap_selector(page, args.selector, args.wait)
    else:
        fw = await page.evaluate(JS_DETECT_FRAMEWORK)
        if fw:
            print(f"  Strategy: {fw['fw']} ({fw['total']} slides)")
            imgs = await cap_framework(page, fw, args.wait)
        else:
            dom = await page.evaluate(JS_DOM_COUNT)
            hint = dom["n"] if dom else None
            print(f"  Strategy: keyboard nav" + (f" (~{hint} slides)" if hint else ""))
            imgs = await cap_keyboard(page, args.wait, hint)
            if len(imgs) <= 1:
                print("  Fallback: scroll chunking")
                await page.goto(url, wait_until="networkidle")
                await page.wait_for_timeout(600)
                imgs = await cap_scroll(page, args.wait)

    await browser.close()
    return imgs


# ── Video recording ──────────────────────────────────────────────────

async def record_mp4(pw, url, output: Path, width: int, height: int,
                     slide_count: int, duration: float, zoom: int = 100):
    """Record the presentation live — all CSS animations captured."""
    tmpdir = tempfile.mkdtemp(prefix="html2deck_")

    browser = await pw.chromium.launch()
    context = await browser.new_context(
        viewport={"width": width, "height": height},
        record_video_dir=tmpdir,
        record_video_size={"width": width, "height": height},
    )
    page = await context.new_page()

    print("  Recording video (this takes real time)…")
    await page.goto(url, wait_until="networkidle", timeout=30_000)
    await page.evaluate(JS_FONTS)
    await page.wait_for_timeout(2000)
    if zoom != 100:
        await page.evaluate(JS_ZOOM, zoom / 100)
        await page.wait_for_timeout(500)

    ms = int(duration * 1000)
    for i in range(slide_count):
        print(f"\r  Recording slide {i+1}/{slide_count} ({duration}s each)", end="", flush=True)
        await page.wait_for_timeout(ms)
        if i < slide_count - 1:
            await page.keyboard.press("ArrowRight")

    # hold last slide
    await page.wait_for_timeout(2000)
    print()

    video_path = await page.video.path()
    await context.close()
    await browser.close()

    # WebM → MP4
    print("  Encoding MP4 (ffmpeg)…")
    r = subprocess.run([
        "ffmpeg", "-y", "-i", str(video_path),
        "-ss", "1.5",                    # trim loading
        "-c:v", "libx264", "-preset", "medium", "-crf", "20",
        "-pix_fmt", "yuv420p",
        "-movflags", "+faststart",
        str(output),
    ], capture_output=True, text=True)

    # cleanup
    try:
        Path(video_path).unlink(missing_ok=True)
        for f in Path(tmpdir).iterdir(): f.unlink(missing_ok=True)
        Path(tmpdir).rmdir()
    except OSError:
        pass

    if r.returncode != 0:
        print(f"  ffmpeg failed: {r.stderr[-300:]}", file=sys.stderr)
        return False
    return True


# ── Assemblers ───────────────────────────────────────────────────────

def assemble_pdf(images, output):
    pil = [Image.open(io.BytesIO(b)).convert("RGB") for b in images]
    pil[0].save(str(output), "PDF", save_all=True, append_images=pil[1:], resolution=150)

def assemble_pptx(images, output, w, h):
    prs = Presentation()
    a = w / h; sh = 7.5; sw = sh * a
    prs.slide_width = Inches(sw); prs.slide_height = Inches(sh)
    blank = min(prs.slide_layouts, key=lambda l: len(list(l.placeholders)))
    for raw in images:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(io.BytesIO(raw), 0, 0, prs.slide_width, prs.slide_height)
    prs.save(str(output))


# ── Main ─────────────────────────────────────────────────────────────

async def run(args):
    html_path = Path(args.input).resolve()
    if not html_path.is_file():
        print(f"Error: {html_path} not found", file=sys.stderr); sys.exit(1)

    stem = args.output or str(html_path.with_suffix(""))

    fmt = args.format
    if fmt == "both": fmts = ["pdf", "pptx"]
    elif fmt == "all": fmts = ["pdf", "pptx", "mp4"]
    else: fmts = [fmt]

    need_ss = "pdf" in fmts or "pptx" in fmts
    need_vid = "mp4" in fmts

    print("html2deck")
    print(f"  Input:  {html_path.name}")
    print(f"  Output: {', '.join(f'{Path(stem).name}.{f}' for f in fmts)}")
    print()

    srv, port = _serve(html_path.parent)
    url = f"http://127.0.0.1:{port}/{html_path.name}"

    try:
        async with async_playwright() as pw:
            images = []
            slide_count = 0

            # Screenshot capture for PDF/PPTX
            if need_ss:
                print("  Capturing screenshots…")
                images = await capture_slides(pw, url, args)
                slide_count = len(images)
                if not images:
                    print("  No slides captured!", file=sys.stderr); sys.exit(1)
                print(f"  {slide_count} slides captured\n")

                if args.debug:
                    dd = Path(f"{stem}_slides"); dd.mkdir(exist_ok=True)
                    for i, p in enumerate(images): (dd / f"slide_{i+1:03d}.png").write_bytes(p)
                    print(f"  Debug PNGs: {dd}/")

                for f in fmts:
                    if f == "mp4": continue
                    out = Path(f"{stem}.{f}")
                    if f == "pdf": assemble_pdf(images, out)
                    else: assemble_pptx(images, out, args.width, args.height)
                    print(f"  {f.upper()} -> {out}  ({out.stat().st_size/1048576:.1f} MB)")

            # Video recording for MP4
            if need_vid:
                if slide_count == 0:
                    # Detect slide count without full capture
                    print("  Detecting slides…")
                    browser = await pw.chromium.launch()
                    page = await browser.new_page(viewport={"width": args.width, "height": args.height})
                    await page.goto(url, wait_until="networkidle", timeout=30_000)
                    await page.evaluate(JS_FONTS)
                    await page.wait_for_timeout(1200)
                    fw = await page.evaluate(JS_DETECT_FRAMEWORK)
                    if fw:
                        slide_count = fw["total"]
                    else:
                        dom = await page.evaluate(JS_DOM_COUNT)
                        slide_count = dom["n"] if dom else 30
                    await browser.close()
                    print(f"  {slide_count} slides detected\n")

                print()
                out = Path(f"{stem}.mp4")
                ok = await record_mp4(pw, url, out, args.width, args.height,
                                      slide_count, args.duration, args.zoom)
                if ok:
                    print(f"  MP4 -> {out}  ({out.stat().st_size/1048576:.1f} MB)")

    finally:
        srv.shutdown()

    print("\n  Done!")


def main():
    p = argparse.ArgumentParser(prog="html2deck",
        description="Convert HTML presentations to PDF, PPTX, and/or MP4",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
examples:
  html2deck slides.html                    # PDF + PPTX (screenshots)
  html2deck slides.html -f mp4             # MP4 video with live effects
  html2deck slides.html -f all             # PDF + PPTX + MP4
  html2deck slides.html -f mp4 --duration 5
  html2deck slides.html --selector ".slide"
        """)
    p.add_argument("input", help="HTML file path")
    p.add_argument("-o", "--output", help="Output base path (no extension)")
    p.add_argument("-f", "--format", choices=["pdf","pptx","mp4","both","all"], default="both",
                   help="Output format (default: both=pdf+pptx)")
    p.add_argument("--width", type=int, default=1920)
    p.add_argument("--height", type=int, default=1080)
    p.add_argument("--scale", type=float, default=2.0, help="Screenshot DPR (default: 2)")
    p.add_argument("--zoom", type=int, default=100,
                   help="CSS zoom %% — 130 = 30%% bigger content (default: 100)")
    p.add_argument("--wait", type=int, default=800, help="Wait ms between slides (default: 800)")
    p.add_argument("--duration", type=float, default=4.0,
                   help="Seconds per slide in MP4 (default: 4)")
    p.add_argument("--selector", help="Force CSS selector for slides")
    p.add_argument("--debug", action="store_true", help="Save individual PNGs")
    args = p.parse_args()
    asyncio.run(run(args))

if __name__ == "__main__":
    main()
